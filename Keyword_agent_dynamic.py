"""
Keyword Intelligence Agent - Dynamic Version
User types keywords when the script runs.
No need to edit the code for different searches.
"""

import os
import re
import fitz
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document


# ─────────────────────────────────────────────
# SECTION 1 — ASK USER FOR INPUT
# ─────────────────────────────────────────────
def get_user_inputs():
    print("\n" + "="*50)
    print("   KEYWORD INTELLIGENCE AGENT")
    print("="*50)

    # ask for folder path with existence check
    while True:
        folder = input("\n Enter folder path where your PDFs are:\n > ").strip()
        if os.path.isdir(folder):
            break
        print(f"  ERROR: Folder not found -> {folder}")
        print("  Please check the path and try again.")

    # ask for must-have keywords
    print("\n Enter MUST-HAVE keywords (all must appear in sentence)")
    print(" Separate multiple keywords with a comma")
    print(" Example:  capital expenditure, revenue")
    must_have_input = input(" > ").strip()

    # ask for any-of keywords
    print("\n Enter ANY-OF keywords (at least one must appear)")
    print(" Separate multiple keywords with a comma")
    print(" Example:  AI, artificial intelligence, machine learning")
    any_of_input = input(" > ").strip()

    # convert comma-separated strings into lists
    must_have = [kw.strip() for kw in must_have_input.split(",") if kw.strip()]
    any_of    = [kw.strip() for kw in any_of_input.split(",")    if kw.strip()]

    # show user what was understood
    print("\n" + "-"*50)
    print(f"  Folder    : {folder}")
    print(f"  Must-have : {must_have}")
    print(f"  Any-of    : {any_of}")
    print("-"*50)

    # confirm before running
    confirm = input("\n Start scanning? (yes / no): ").strip().lower()
    if confirm != "yes":
        print("\n  Cancelled. Run the script again when ready.")
        exit()

    return folder, must_have, any_of


# ─────────────────────────────────────────────
# SECTION 2 — KEYWORD LOGIC
# ─────────────────────────────────────────────
def sentence_matches(sentence, must_have, any_of):
    s = sentence.lower()
    return all(kw.lower() in s for kw in must_have) \
       and any(kw.lower() in s for kw in any_of)


# ─────────────────────────────────────────────
# SECTION 3 — SENTENCE SPLITTER
# ─────────────────────────────────────────────
def split_into_sentences(text):
    sentences = re.split(r'(?<=[.!?])\s+', text)
    return [s.replace("\n", " ").strip() for s in sentences if s.strip()]


# ─────────────────────────────────────────────
# SECTION 4 — READ PDF
# ─────────────────────────────────────────────
def extract_sentences_from_pdf(pdf_path):
    results = []
    doc = fitz.open(pdf_path)
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text")
        sentences = split_into_sentences(text)
        for sentence in sentences:
            if sentence.strip():
                results.append((page_num, sentence.strip()))
    doc.close()
    return results


# ─────────────────────────────────────────────
# SECTION 5 — READ DOCX
# ─────────────────────────────────────────────
def extract_sentences_from_docx(docx_path):
    results = []
    doc = Document(docx_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    page_size = 30
    for idx, para in enumerate(paragraphs):
        page_num = (idx // page_size) + 1
        sentences = split_into_sentences(para)
        for sentence in sentences:
            if sentence.strip():
                results.append((page_num, sentence.strip()))
    return results


# ─────────────────────────────────────────────
# SECTION 6 — SEARCH ALL FILES
# ─────────────────────────────────────────────
def search_all_files(folder, must_have, any_of):
    matches = []
    for filename in sorted(os.listdir(folder)):
        filepath = os.path.join(folder, filename)
        ext = filename.lower().split(".")[-1]

        if ext == "pdf":
            print(f"  Scanning PDF : {filename}")
            sentences = extract_sentences_from_pdf(filepath)
        elif ext == "docx":
            print(f"  Scanning DOCX: {filename}")
            sentences = extract_sentences_from_docx(filepath)
        else:
            continue

        for page_num, sentence in sentences:
            if sentence_matches(sentence, must_have, any_of):
                matches.append({
                    "File Name":        filename,
                    "Page Number":      page_num,
                    "Matched Sentence": sentence,
                })
    return matches


# ─────────────────────────────────────────────
# SECTION 7 — SAVE TO EXCEL
# ─────────────────────────────────────────────
def save_to_excel(matches, folder):
    out_path = os.path.join(folder, "Keyword_Research_Results.xlsx")
    # if file is open in Excel, warn and use a timestamped name instead of crashing
    if os.path.exists(out_path):
        try:
            os.rename(out_path, out_path)
        except PermissionError:
            from datetime import datetime
            stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            out_path = os.path.join(folder, f"Keyword_Research_Results_{stamp}.xlsx")
            print(f"  WARNING: Previous Excel file is open. Saving as: {os.path.basename(out_path)}")
    df = pd.DataFrame(matches, columns=["File Name", "Page Number", "Matched Sentence"])
    with pd.ExcelWriter(out_path, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Matches")
        ws = writer.sheets["Matches"]
        for col in ws.columns:
            max_len = max((len(str(cell.value)) for cell in col if cell.value), default=10)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 80)
    print(f"  Excel saved  : {out_path}")
    return out_path


# ─────────────────────────────────────────────
# SECTION 8 — SAVE TO PPTX
# ─────────────────────────────────────────────
def save_to_pptx(matches, folder):
    out_path = os.path.join(folder, "Keyword_Research_Results.pptx")
    if os.path.exists(out_path):
        try:
            os.rename(out_path, out_path)
        except PermissionError:
            from datetime import datetime
            stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            out_path = os.path.join(folder, f"Keyword_Research_Results_{stamp}.pptx")
            print(f"  WARNING: Previous PPTX file is open. Saving as: {os.path.basename(out_path)}")
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for match in matches:
        slide = prs.slides.add_slide(blank_layout)

        # header bar
        header = slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(13.33), Inches(1.1))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)
        tf = header.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text           = match["File Name"]
        run.font.size      = Pt(20)
        run.font.bold      = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # matched sentence
        body = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.0))
        tf2 = body.text_frame
        tf2.word_wrap = True
        run2 = tf2.paragraphs[0].add_run()
        run2.text           = match["Matched Sentence"]
        run2.font.size      = Pt(16)
        run2.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

        # page number label
        label = slide.shapes.add_textbox(
            Inches(10.5), Inches(6.7), Inches(2.5), Inches(0.5))
        tf3 = label.text_frame
        run3 = tf3.paragraphs[0].add_run()
        run3.text           = f"Page {match['Page Number']}"
        run3.font.size      = Pt(11)
        run3.font.italic    = True
        run3.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    prs.save(out_path)
    print(f"  PPTX saved   : {out_path}")
    return out_path


# ─────────────────────────────────────────────
# SECTION 9 — MAIN
# ─────────────────────────────────────────────
def main():

    # step 1 — ask user for all inputs
    folder, must_have, any_of = get_user_inputs()

    # step 2 — scan all files
    print("\n[ Scanning files... ]\n")
    matches = search_all_files(folder, must_have, any_of)

    # step 3 — check if anything matched
    if not matches:
        print("\n  No sentences matched. Try different keywords.")
        return

    print(f"\n  {len(matches)} match(es) found.\n")

    # step 4 — save outputs
    print("[ Saving Excel... ]")
    save_to_excel(matches, folder)

    print("[ Building PowerPoint... ]")
    save_to_pptx(matches, folder)

    print("\n" + "="*50)
    print("   DONE! Both files saved to your folder.")
    print("="*50 + "\n")


if __name__ == "__main__":
    main()
